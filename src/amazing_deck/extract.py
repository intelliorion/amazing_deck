"""Template analysis — extract layouts, theme, images, and write knowledge hub."""
from pathlib import Path
from xml.etree import ElementTree as ET
import json

from pptx import Presentation

from .utils import emu_to_inches, slugify
from .thumbnails import render_layout_thumbnails


NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def analyze_template(template_path, output_dir):
    """Main entry point — template → knowledge hub directory."""
    template_path = Path(template_path)
    hub = Path(output_dir)
    hub.mkdir(parents=True, exist_ok=True)
    (hub / "layouts").mkdir(exist_ok=True)
    (hub / "assets").mkdir(exist_ok=True)
    (hub / "assets" / "images").mkdir(exist_ok=True)

    prs = Presentation(str(template_path))

    theme = extract_theme(prs)
    backgrounds = extract_backgrounds(prs, theme)
    layouts = [describe_layout(i, layout)
               for i, layout in enumerate(prs.slide_layouts)]
    image_count = extract_images(prs, hub / "assets" / "images")
    thumbnails = render_layout_thumbnails(
        template_path, prs, hub / "layouts" / "thumbnails", backgrounds)

    bg_by_idx = {b["index"]: b for b in backgrounds.get("layouts", [])}
    for lm in layouts:
        lm["thumbnail"] = thumbnails.get(lm["index"])
        lm["background"] = bg_by_idx.get(lm["index"])

    manifest = {
        "template_path": str(template_path),
        "template_name": template_path.stem,
        "slide_size_in": {
            "width": emu_to_inches(prs.slide_width),
            "height": emu_to_inches(prs.slide_height),
            "aspect": _aspect_ratio(prs.slide_width, prs.slide_height),
        },
        "layouts": layouts,
        "theme": theme,
        "backgrounds": backgrounds,
        "image_count": image_count,
    }

    (hub / "manifest.json").write_text(json.dumps(manifest, indent=2))
    (hub / "assets" / "backgrounds.json").write_text(json.dumps(backgrounds, indent=2))
    _write_overview(hub, manifest)
    _write_style_guide(hub, theme, backgrounds)
    _write_colors_json(hub, theme)
    _write_fonts_json(hub, theme)
    for lm in layouts:
        _write_layout_md(hub / "layouts", lm)
    _write_generation_prompt(hub, manifest)

    return hub


def describe_layout(idx, layout):
    """Return a dict describing a single layout."""
    placeholders = []
    for ph in layout.placeholders:
        pf = ph.placeholder_format
        ph_type = _norm_type(pf.type) if pf.type else "UNKNOWN"
        placeholders.append({
            "idx": pf.idx,
            "type": ph_type,
            "role": _type_to_role(ph_type),
            "name": ph.name,
            "has_text_frame": ph.has_text_frame,
            "default_text": (ph.text_frame.text if ph.has_text_frame else "").strip(),
            "left_in": emu_to_inches(ph.left),
            "top_in": emu_to_inches(ph.top),
            "width_in": emu_to_inches(ph.width),
            "height_in": emu_to_inches(ph.height),
        })
    return {
        "index": idx,
        "name": layout.name,
        "slug": slugify(layout.name),
        "placeholders": placeholders,
    }


def extract_theme(prs):
    """Pull theme colors and fonts from the slide master's theme XML."""
    colors = {}
    fonts = {}
    try:
        master_part = prs.slide_master.part
        for rel in master_part.rels.values():
            if "theme" in rel.reltype.lower():
                theme_xml = rel.target_part.blob.decode("utf-8")
                root = ET.fromstring(theme_xml)
                for scheme in root.iter(f"{{{NS_A}}}clrScheme"):
                    for c in scheme:
                        tag = c.tag.rsplit("}", 1)[-1]
                        srgb = c.find(f"{{{NS_A}}}srgbClr")
                        sysc = c.find(f"{{{NS_A}}}sysClr")
                        if srgb is not None:
                            colors[tag] = "#" + srgb.get("val", "000000").upper()
                        elif sysc is not None:
                            colors[tag] = "#" + sysc.get("lastClr", "000000").upper()
                for scheme in root.iter(f"{{{NS_A}}}fontScheme"):
                    major = scheme.find(f"{{{NS_A}}}majorFont/{{{NS_A}}}latin")
                    minor = scheme.find(f"{{{NS_A}}}minorFont/{{{NS_A}}}latin")
                    if major is not None:
                        fonts["major"] = major.get("typeface")
                    if minor is not None:
                        fonts["minor"] = minor.get("typeface")
                break
    except Exception as exc:
        print(f"  [warn] Could not extract theme: {exc}")
    return {"colors": colors, "fonts": fonts}


def extract_backgrounds(prs, theme):
    """Extract <p:bg> from master + each layout. Resolve theme colors to hex."""
    ns = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": NS_A,
    }
    colors = theme.get("colors", {})

    def parse_bg(element):
        bg = element.find(".//p:bg", ns)
        if bg is None:
            return None
        bg_pr = bg.find("p:bgPr", ns)
        if bg_pr is None:
            return {"fill_type": "inherit", "hex": None}

        solid = bg_pr.find("a:solidFill", ns)
        if solid is not None:
            scheme = solid.find("a:schemeClr", ns)
            srgb = solid.find("a:srgbClr", ns)
            if scheme is not None:
                scheme_name = scheme.get("val", "")
                aliases = {"tx1": "dk1", "bg1": "lt1", "tx2": "dk2", "bg2": "lt2"}
                resolved = aliases.get(scheme_name, scheme_name)
                hex_val = colors.get(resolved, "#CCCCCC")
                lm = scheme.find("a:lumMod", ns)
                lo = scheme.find("a:lumOff", ns)
                lm_val = int(lm.get("val")) if lm is not None else None
                lo_val = int(lo.get("val")) if lo is not None else None
                if lm_val is not None or lo_val is not None:
                    hex_val = apply_luminance(hex_val, lm_val, lo_val)
                return {
                    "fill_type": "solid",
                    "hex": hex_val,
                    "theme_color": scheme_name,
                    "lumMod": lm_val,
                    "lumOff": lo_val,
                }
            if srgb is not None:
                return {"fill_type": "solid", "hex": "#" + srgb.get("val", "000000").upper()}

        if bg_pr.find("a:gradFill", ns) is not None:
            return {"fill_type": "gradient", "hex": "#CCCCCC"}
        if bg_pr.find("a:blipFill", ns) is not None:
            return {"fill_type": "picture", "hex": "#CCCCCC"}
        if bg_pr.find("a:pattFill", ns) is not None:
            return {"fill_type": "pattern", "hex": "#CCCCCC"}

        return {"fill_type": "unknown", "hex": "#CCCCCC"}

    master_bg = parse_bg(prs.slide_master.element) or {"fill_type": "default", "hex": "#FFFFFF"}

    layout_bgs = []
    for i, layout in enumerate(prs.slide_layouts):
        bg = parse_bg(layout.element)
        if bg is None:
            bg = {**master_bg, "inherited": True}
        else:
            bg["inherited"] = False
        bg["index"] = i
        bg["name"] = layout.name
        layout_bgs.append(bg)

    return {"master": master_bg, "layouts": layout_bgs}


def apply_luminance(hex_color, lum_mod=None, lum_off=None):
    """Apply OOXML lumMod/lumOff modifiers. Values are 1/1000 units (100000 = 100%)."""
    import colorsys
    s = str(hex_color).lstrip("#")
    if len(s) != 6:
        return hex_color
    try:
        r = int(s[0:2], 16) / 255
        g = int(s[2:4], 16) / 255
        b = int(s[4:6], 16) / 255
    except ValueError:
        return hex_color
    h, l, sat = colorsys.rgb_to_hls(r, g, b)
    if lum_mod is not None:
        l = l * (lum_mod / 100000.0)
    if lum_off is not None:
        l = l + (lum_off / 100000.0)
    l = max(0.0, min(1.0, l))
    r, g, b = colorsys.hls_to_rgb(h, l, sat)
    return "#{:02X}{:02X}{:02X}".format(int(r * 255), int(g * 255), int(b * 255))


def extract_images(prs, out_dir):
    """Extract every image embedded in slide masters + layouts."""
    count = 0
    try:
        for m_idx, master in enumerate(prs.slide_masters):
            count += _extract_from_shapes(master.shapes, out_dir,
                                          f"master{m_idx}")
            for layout in master.slide_layouts:
                count += _extract_from_shapes(
                    layout.shapes, out_dir,
                    f"layout_{slugify(layout.name)}")
    except Exception as exc:
        print(f"  [warn] Image extraction issue: {exc}")
    return count


def _extract_from_shapes(shapes, out_dir, prefix):
    count = 0
    for shape in shapes:
        try:
            if shape.shape_type == 13:  # PICTURE
                img = shape.image
                name = f"{prefix}_{shape.shape_id}.{img.ext}"
                (out_dir / name).write_bytes(img.blob)
                count += 1
        except Exception:
            continue
    return count


def _norm_type(t):
    """Normalize PP_PLACEHOLDER enum string across python-pptx versions.

    Handles:
      'CENTER_TITLE (3)'          -> 'CENTER_TITLE'
      'PP_PLACEHOLDER.TITLE'      -> 'TITLE'
      'TITLE'                     -> 'TITLE'
    """
    if t is None:
        return "UNKNOWN"
    try:
        name = getattr(t, "name", None)
        if name:
            return name
    except Exception:
        pass
    s = str(t)
    s = s.split(" ", 1)[0]
    s = s.rsplit(".", 1)[-1]
    return s or "UNKNOWN"


def _type_to_role(ph_type):
    mapping = {
        "TITLE": "title",
        "CENTER_TITLE": "title",
        "CTR_TITLE": "title",
        "SUBTITLE": "subtitle",
        "BODY": "content",
        "OBJECT": "content",
        "DATE": "date",
        "FOOTER": "footer",
        "SLIDE_NUMBER": "slide_number",
        "PICTURE": "image",
    }
    return mapping.get(ph_type, ph_type.lower())


def _aspect_ratio(w, h):
    if not w or not h:
        return "unknown"
    ratio = w / h
    if abs(ratio - 16 / 9) < 0.01:
        return "16:9"
    if abs(ratio - 4 / 3) < 0.01:
        return "4:3"
    if abs(ratio - 16 / 10) < 0.01:
        return "16:10"
    return f"{w}x{h}"


# ---------- Writers ----------

def _write_overview(hub, manifest):
    w = manifest["slide_size_in"]["width"]
    h = manifest["slide_size_in"]["height"]
    aspect = manifest["slide_size_in"]["aspect"]
    n_layouts = len(manifest["layouts"])
    n_images = manifest["image_count"]
    theme_colors = manifest["theme"]["colors"]
    theme_fonts = manifest["theme"]["fonts"]
    layouts = manifest["layouts"]

    lines = [
        f"# Overview — {manifest['template_name']}",
        "",
        f"- **Source:** `{manifest['template_path']}`",
        f"- **Dimensions:** {w} x {h} in  ({aspect})",
        f"- **Layouts:** {n_layouts}",
        f"- **Theme colors:** {len(theme_colors)}",
        f"- **Theme fonts:** major={theme_fonts.get('major', 'n/a')}, "
        f"minor={theme_fonts.get('minor', 'n/a')}",
        f"- **Embedded images:** {n_images} (see `assets/images/`)",
        "",
        "## Layouts at a glance",
        "",
    ]
    cols = 3
    rows = (len(layouts) + cols - 1) // cols
    header = "| " + " | ".join([" "] * cols) + " |"
    sep = "|" + "|".join(["---"] * cols) + "|"
    lines += [header, sep]
    for r in range(rows):
        imgs, labels = [], []
        for c in range(cols):
            idx = r * cols + c
            if idx < len(layouts):
                lm = layouts[idx]
                thumb = lm.get("thumbnail") or ""
                imgs.append(f"![{lm['index']}](layouts/{thumb})" if thumb else " ")
                labels.append(f"**[{lm['index']}] {lm['name']}**")
            else:
                imgs.append(" ")
                labels.append(" ")
        lines.append("| " + " | ".join(imgs) + " |")
        lines.append("| " + " | ".join(labels) + " |")
    lines += ["", "## Layouts — detailed list", ""]
    for lm in layouts:
        roles = ", ".join(sorted({p["role"] for p in lm["placeholders"]})) or "none"
        bg = (lm.get("background") or {}).get("hex") or "inherit"
        lines.append(f"- **[{lm['index']}] {lm['name']}** — "
                     f"{len(lm['placeholders'])} placeholders ({roles}) · bg `{bg}`")
    lines += ["",
              "See `layouts/` for per-layout detail, "
              "`style-guide.md` for colors/fonts, and "
              "`assets/backgrounds.json` for machine-readable background data."]
    (hub / "overview.md").write_text("\n".join(lines))


def _write_style_guide(hub, theme, backgrounds=None):
    colors = theme.get("colors", {})
    fonts = theme.get("fonts", {})
    lines = [
        "# Style Guide",
        "",
        "## Theme colors",
        "",
        "| Role | Hex |",
        "|---|---|",
    ]
    for role, hex_val in colors.items():
        lines.append(f"| `{role}` | `{hex_val}` |")
    lines += [
        "",
        "## Fonts",
        "",
        f"- **Major (headings):** `{fonts.get('major', 'n/a')}`",
        f"- **Minor (body):** `{fonts.get('minor', 'n/a')}`",
        "",
    ]
    if backgrounds:
        lines += ["## Backgrounds by layout", "", "| Layout | Fill type | Color | Theme source |", "|---|---|---|---|"]
        master = backgrounds.get("master", {})
        mfill = master.get("fill_type", "?")
        mhex = master.get("hex") or "—"
        lines.append(f"| **Master** | {mfill} | `{mhex}` | — |")
        for b in backgrounds.get("layouts", []):
            inherited = " (inherits)" if b.get("inherited") else ""
            src = b.get("theme_color") or "—"
            lines.append(f"| [{b['index']}] {b['name']}{inherited} | "
                         f"{b.get('fill_type','?')} | `{b.get('hex') or '—'}` | `{src}` |")
        lines.append("")
    lines += [
        "## Generation rules",
        "",
        "When producing content for this template:",
        "- Prefer theme colors (above) over arbitrary hex values.",
        "- Use the major font for titles, the minor font for body copy.",
        "- Body text minimum 10pt; 14pt preferred for projected decks.",
        "- Maximum 7 bullets per slide. Split instead of cramming.",
        "- One message per slide. Two messages = two slides.",
        "- Use the background colors above to pick contrasting text (dark bg = light text).",
    ]
    (hub / "style-guide.md").write_text("\n".join(lines))


def _write_colors_json(hub, theme):
    (hub / "assets" / "colors.json").write_text(
        json.dumps(theme.get("colors", {}), indent=2))


def _write_fonts_json(hub, theme):
    (hub / "assets" / "fonts.json").write_text(
        json.dumps(theme.get("fonts", {}), indent=2))


def _write_layout_md(layouts_dir, lm):
    filename = f"{lm['index']:02d}-{lm['slug']}.md"
    bg = lm.get("background") or {}
    bg_line = ""
    if bg.get("hex"):
        fill = bg.get("fill_type", "")
        source = f" ({bg.get('theme_color')})" if bg.get("theme_color") else ""
        inherit = " — inherits from master" if bg.get("inherited") else ""
        bg_line = f"- **Background:** {fill} `{bg['hex']}`{source}{inherit}"
    thumb = lm.get("thumbnail")
    thumb_md = f"![Thumbnail]({thumb.split('/')[-1] if thumb and '/' in thumb else thumb})" if thumb else ""
    # thumb path is like "thumbnails/NN-name.png" relative to layouts/
    thumb_md = f"![Thumbnail]({thumb})" if thumb else ""
    lines = [
        f"# Layout {lm['index']}: {lm['name']}",
        "",
        thumb_md,
        "" if not thumb_md else "",
        f"- **Index:** {lm['index']}",
        f"- **Slug:** `{lm['slug']}`",
        f"- **Placeholders:** {len(lm['placeholders'])}",
    ]
    if bg_line:
        lines.append(bg_line)
    lines += [
        "",
        "## Placeholders",
        "",
        "| idx | role | name | size (in) | default text |",
        "|---|---|---|---|---|",
    ]
    for p in lm["placeholders"]:
        size = f"{p['width_in']}x{p['height_in']}" if p["width_in"] else "?"
        dt = p["default_text"][:40].replace("|", "\\|") or "_(empty)_"
        lines.append(f"| {p['idx']} | {p['role']} | `{p['name']}` | {size} | {dt} |")
    lines += [
        "",
        "## Example content JSON",
        "",
        "```json",
        "{",
        f'  "layout": "{lm["name"]}",',
        '  "placeholders": {',
    ]
    keys = []
    for p in lm["placeholders"]:
        sample = _sample_for_role(p["role"])
        keys.append(f'    "{p["role"]}": {sample}')
    lines.append(",\n".join(keys))
    lines += ["  }", "}", "```"]
    (layouts_dir / filename).write_text("\n".join(lines))


def _sample_for_role(role):
    if role == "content":
        return '["First point", "Second point", "Third point"]'
    if role == "title":
        return '"Slide Title"'
    if role == "subtitle":
        return '"Supporting subtitle"'
    if role == "date":
        return '"2026-Q3"'
    return '"..."'


def _write_generation_prompt(hub, manifest):
    """The killer artifact - a prompt LLMs can use to generate on-brand content."""
    colors = manifest["theme"]["colors"]
    fonts = manifest["theme"]["fonts"]
    layouts = manifest["layouts"]

    lines = [
        f"# Generation Prompt - {manifest['template_name']}",
        "",
        "Paste this prompt (plus your content outline) into any LLM to produce",
        "a content.json compatible with amazing-deck.",
        "",
        "---",
        "",
        "## TEMPLATE RULES",
        "",
        f"You are generating slide content for the **{manifest['template_name']}** template.",
        "",
        f"- **Dimensions:** {manifest['slide_size_in']['aspect']}",
        f"- **Major font:** {fonts.get('major', 'n/a')} (titles)",
        f"- **Minor font:** {fonts.get('minor', 'n/a')} (body)",
        "",
        "**Color palette (use these, not arbitrary hex):**",
    ]
    for role, hex_val in colors.items():
        lines.append(f"- `{role}`: `{hex_val}`")

    lines += [
        "",
        "## AVAILABLE LAYOUTS",
        "",
        "Use these layout names (exact match):",
        "",
    ]
    for lm in layouts:
        roles = ", ".join(sorted({p["role"] for p in lm["placeholders"]})) or "none"
        lines.append(f"- `{lm['name']}` - placeholders: {roles}")

    lines += [
        "",
        "## AVAILABLE RECIPES",
        "",
        "For visual patterns not in the template, use recipes:",
        "",
        "- `kpi_cards` - dashboards (3-4 big numbers)",
        "- `comparison` - 2-column do/don't or before/after",
        "- `timeline` - milestones with dates and status",
        "- `chart_bar` - native PowerPoint bar chart",
        "- `asks` - numbered executive asks",
        "",
        "## QUALITY RULES (non-negotiable)",
        "",
        "1. Max 7 bullets per slide.",
        "2. One message per slide - if two, split.",
        "3. Use layout names exactly as listed above.",
        "4. Use recipes for charts, timelines, KPIs - never fake them with textboxes.",
        "5. Open with a commitment slide (headline metric, target, date).",
        "6. Close with an asks slide using the `asks` recipe.",
        "",
        "## OUTPUT FORMAT",
        "",
        "Return a valid JSON object with this shape:",
        "",
        "```json",
        "{",
        '  "slides": [',
        '    {"layout": "LayoutName", "placeholders": {"role": "value"}},',
        '    {"recipe": "recipe_name", "recipe_params": {}}',
        "  ]",
        "}",
        "```",
        "",
        "---",
        "",
        "## YOUR CONTENT (user fills in below)",
        "",
        "_Describe what you want in the deck - title, audience, key messages,",
        "data to chart, asks to leadership, etc._",
    ]
    (hub / "generation-prompt.md").write_text("\n".join(lines))
