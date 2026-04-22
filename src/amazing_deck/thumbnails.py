"""Layout thumbnail rendering.

Two rendering paths:
  - High fidelity: shell out to LibreOffice (`soffice --headless --convert-to png`)
  - Fallback: PIL schematic (background color + placeholder outlines)

Selected per-call based on whether `soffice` is on PATH. Each layout that
fails the high-fidelity path individually falls back to schematic, without
aborting the whole analyze run.
"""
from pathlib import Path
import shutil
import subprocess
import tempfile

from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont

from .utils import slugify


# ----------------------------- public API -----------------------------

def render_layout_thumbnails(template_path, prs, output_dir, backgrounds):
    """Render a PNG thumbnail per layout.

    Args:
        template_path: Path to the source .pptx (needed for soffice path)
        prs: loaded Presentation object
        output_dir: Path — where PNGs will be written
        backgrounds: dict with 'layouts' list of {index, hex, fill_type, ...}

    Returns:
        dict[int, str]: {layout_index: relative_png_path_from_output_dir_parent}
    """
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    use_soffice = _has_soffice()
    mode = "LibreOffice (high fidelity)" if use_soffice else "PIL schematic (no soffice found)"
    print(f"  Thumbnails: {mode}")

    bg_by_idx = {b["index"]: b for b in backgrounds.get("layouts", [])}
    result = {}

    for i, layout in enumerate(prs.slide_layouts):
        slug = slugify(layout.name)
        out_png = output_dir / f"{i:02d}-{slug}.png"

        rendered = False
        if use_soffice:
            try:
                rendered = _render_with_soffice(template_path, i, layout, out_png)
            except Exception as exc:
                print(f"  [warn] soffice exception on layout {i} '{layout.name}': {exc}")

        if not rendered:
            _render_schematic(layout, bg_by_idx.get(i, {}), out_png)

        result[i] = f"thumbnails/{out_png.name}"

    return result


# ----------------------------- soffice -----------------------------

def _has_soffice():
    return shutil.which("soffice") is not None or shutil.which("soffice.com") is not None


def _render_with_soffice(template_path, layout_idx, layout, out_png):
    """Build a minimal one-slide pptx using this layout, convert with soffice."""
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        mini_pptx = tmp_path / f"layout_{layout_idx}.pptx"

        # Start from the source template to keep theme/master intact
        prs = Presentation(str(template_path))

        # Clear all existing slides (drop the sldId elements)
        from pptx.oxml.ns import qn
        xml_slides = prs.slides._sldIdLst
        for sld_id in list(xml_slides):
            r_id = sld_id.get(qn("r:id"))
            if r_id:
                try:
                    prs.part.drop_rel(r_id)
                except Exception:
                    pass
            xml_slides.remove(sld_id)

        # Add exactly one empty slide using the requested layout
        prs.slides.add_slide(list(prs.slide_layouts)[layout_idx])
        prs.save(str(mini_pptx))

        # Shell out to LibreOffice
        try:
            res = subprocess.run(
                ["soffice", "--headless", "--convert-to", "png",
                 "--outdir", str(tmp_path), str(mini_pptx)],
                capture_output=True, timeout=45, text=True,
            )
        except (subprocess.TimeoutExpired, FileNotFoundError):
            return False

        if res.returncode != 0:
            return False

        generated = tmp_path / f"{mini_pptx.stem}.png"
        if not generated.exists():
            return False

        # Resize if excessively wide; keep 16:9
        img = Image.open(generated)
        if img.width > 1600:
            ratio = 1600 / img.width
            img = img.resize((1600, int(img.height * ratio)), Image.LANCZOS)
        img.convert("RGB").save(out_png, "PNG", optimize=True)
        return True


# ----------------------------- PIL schematic -----------------------------

def _render_schematic(layout, bg_info, out_png):
    """Draw a schematic: background color + placeholder wireframes + name strip."""
    W, H = 800, 450  # 16:9

    bg_hex = bg_info.get("hex") or "#CCCCCC"
    bg_rgb = _hex_tuple(bg_hex)
    img = Image.new("RGB", (W, H), bg_rgb)
    draw = ImageDraw.Draw(img)

    dark = _is_dark(bg_hex)
    outline = (240, 240, 240) if dark else (60, 60, 60)
    label_col = (230, 230, 230) if dark else (40, 40, 40)
    strip_col = (40, 40, 40) if dark else (30, 30, 30)
    strip_text = (255, 255, 255)

    font_label = _font(12)
    font_title = _font(14)

    scale_x = W / 13.333
    scale_y = H / 7.5

    for ph in layout.placeholders:
        if ph.left is None or ph.top is None:
            continue
        try:
            x = int(ph.left / 914400 * scale_x)
            y = int(ph.top / 914400 * scale_y)
            w = int((ph.width or 0) / 914400 * scale_x)
            h = int((ph.height or 0) / 914400 * scale_y)
        except Exception:
            continue

        if w <= 0 or h <= 0:
            continue

        x2 = min(x + w, W - 1)
        y2 = min(y + h, H - 1)
        x = max(0, x)
        y = max(0, y)

        draw.rectangle([x, y, x2, y2], outline=outline, width=2)

        role = _ph_role(ph)
        if role:
            draw.text((x + 5, y + 5), role, fill=label_col, font=font_label)

    strip_h = 26
    draw.rectangle([0, H - strip_h, W, H], fill=strip_col)
    draw.text((10, H - strip_h + 5), layout.name, fill=strip_text, font=font_title)
    if bg_info.get("hex"):
        bg_hex_display = bg_info["hex"]
        right_text = f"bg {bg_hex_display}"
        try:
            tw = draw.textlength(right_text, font=font_label)
        except Exception:
            tw = 100
        draw.text((W - tw - 10, H - strip_h + 7), right_text, fill=strip_text, font=font_label)

    img.save(out_png, "PNG", optimize=True)


# ----------------------------- helpers -----------------------------

def _hex_tuple(hex_str):
    s = str(hex_str or "").lstrip("#")
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    if len(s) != 6:
        return (204, 204, 204)
    try:
        return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except ValueError:
        return (204, 204, 204)


def _is_dark(hex_str):
    r, g, b = _hex_tuple(hex_str)
    return (0.299 * r + 0.587 * g + 0.114 * b) < 128


def _ph_role(ph):
    try:
        t = ph.placeholder_format.type
        if t is None:
            return None
        s = str(t).split(" ", 1)[0].rsplit(".", 1)[-1]
        mapping = {
            "TITLE": "title", "CENTER_TITLE": "title", "CTR_TITLE": "title",
            "SUBTITLE": "subtitle", "BODY": "content", "OBJECT": "content",
            "DATE": "date", "FOOTER": "footer", "SLIDE_NUMBER": "page",
            "PICTURE": "image",
        }
        return mapping.get(s, s.lower())
    except Exception:
        return None


def _font(size):
    """Return a PIL font, trying a few common paths."""
    paths = [
        "/System/Library/Fonts/Helvetica.ttc",
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "C:\\Windows\\Fonts\\arial.ttf",
    ]
    for p in paths:
        try:
            return ImageFont.truetype(p, size)
        except (OSError, IOError):
            continue
    return ImageFont.load_default()
