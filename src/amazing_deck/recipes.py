"""Recipe library - compound slide components beyond what templates provide."""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from .render import add_textbox, add_rect
from .utils import hex_to_rgb


REGISTRY = {}


def recipe(name):
    def deco(fn):
        REGISTRY[name] = fn
        return fn
    return deco


def apply_recipe(slide, name, params, theme):
    fn = REGISTRY.get(name)
    if not fn:
        print(f"  [warn] Unknown recipe: {name}. Available: {list(REGISTRY)}")
        return False
    fn(slide, params or {}, theme or {"colors": {}, "fonts": {}})
    return True


def _palette(theme):
    c = (theme or {}).get("colors", {})
    return {
        "primary": c.get("dk2") or c.get("accent1") or "#1B3E6F",
        "accent": c.get("accent1") or "#007DBA",
        "accent2": c.get("accent2") or "#2E8B57",
        "accent3": c.get("accent3") or "#E8A83A",
        "text": c.get("dk1") or "#000000",
        "bg": c.get("lt1") or "#FFFFFF",
        "muted": "#595959",
    }


@recipe("kpi_cards")
def kpi_cards(slide, params, theme):
    """Big-number dashboard.
    Params: {title, cards: [{label, value, subtext?, color?}]}
    """
    p = _palette(theme)
    title = params.get("title")
    cards = params.get("cards", [])

    y = 1.3
    if title:
        add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.5),
                    title, size=18, bold=True, color=p["primary"])
        y = 1.85

    n = len(cards)
    if n == 0:
        return
    total_w = 12.1
    gap = 0.25
    card_w = (total_w - gap * (n - 1)) / n

    for i, card in enumerate(cards):
        x = 0.6 + i * (card_w + gap)
        color = card.get("color") or p["accent"]
        add_rect(slide, Inches(x), Inches(y), Inches(card_w), Inches(1.9),
                 "#FFFFFF", line_hex="#BFBFBF")
        add_rect(slide, Inches(x), Inches(y), Inches(0.12), Inches(1.9), color)
        add_textbox(slide, Inches(x + 0.25), Inches(y + 0.1),
                    Inches(card_w - 0.35), Inches(0.35),
                    card.get("label", ""), size=10, bold=True, color=p["muted"])
        add_textbox(slide, Inches(x + 0.25), Inches(y + 0.45),
                    Inches(card_w - 0.35), Inches(0.95),
                    str(card.get("value", "")), size=32, bold=True,
                    color=color, anchor="middle")
        if card.get("subtext"):
            add_textbox(slide, Inches(x + 0.25), Inches(y + 1.45),
                        Inches(card_w - 0.35), Inches(0.4),
                        card["subtext"], size=9, italic=True, color=p["muted"])


@recipe("comparison")
def comparison(slide, params, theme):
    """Two-column comparison.
    Params: {title, left: {header, items[]}, right: {header, items[]}}
    """
    p = _palette(theme)
    title = params.get("title", "")
    left = params.get("left", {})
    right = params.get("right", {})

    if title:
        add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.5),
                    title, size=18, bold=True, color=p["primary"])

    _col(slide, Inches(0.6), Inches(1.7), Inches(6), Inches(5.2),
         left, params.get("left_color") or p["accent2"], p["text"])
    _col(slide, Inches(6.85), Inches(1.7), Inches(6), Inches(5.2),
         right, params.get("right_color") or p["accent3"], p["text"])


def _col(slide, x, y, w, h, col, header_color, text_color):
    add_rect(slide, x, y, w, Inches(0.7), header_color)
    add_textbox(slide, x, y, w, Inches(0.7),
                col.get("header", ""), size=14, bold=True,
                color="#FFFFFF", align="center", anchor="middle")
    body_y = y + Inches(0.8)
    body_h = h - Inches(0.9)
    add_rect(slide, x, body_y, w, body_h, "#F5F5F5")
    tb = slide.shapes.add_textbox(x + Inches(0.2), body_y + Inches(0.15),
                                    w - Inches(0.4), body_h - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(col.get("items", [])):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.line_spacing = 1.25
        run = par.add_run()
        run.text = f"*  {item}"
        run.font.size = Pt(12)
        run.font.color.rgb = hex_to_rgb(text_color)


@recipe("timeline")
def timeline(slide, params, theme):
    """Horizontal milestones with dates and status.
    Params: {title, milestones: [{date, label, status: done|current|upcoming}]}
    """
    p = _palette(theme)
    title = params.get("title", "")
    ms = params.get("milestones", [])

    if title:
        add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.5),
                    title, size=18, bold=True, color=p["primary"])

    if not ms:
        return

    line_y = 4.0
    x0, x1 = 1.0, 12.3
    add_rect(slide, Inches(x0), Inches(line_y), Inches(x1 - x0),
             Inches(0.08), p["primary"])

    n = len(ms)
    step = (x1 - x0) / max(n - 1, 1) if n > 1 else 0
    status_color = {
        "done": p["accent2"],
        "current": p["accent3"],
        "upcoming": "#999999",
    }

    for i, m in enumerate(ms):
        cx = x0 + step * i if n > 1 else (x0 + x1) / 2
        col = status_color.get(m.get("status", "upcoming"), p["accent"])
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(cx - 0.22), Inches(line_y - 0.18),
            Inches(0.44), Inches(0.44))
        marker.fill.solid()
        marker.fill.fore_color.rgb = hex_to_rgb(col)
        marker.line.color.rgb = hex_to_rgb("#FFFFFF")
        marker.line.width = Pt(2)
        add_textbox(slide, Inches(cx - 1), Inches(line_y - 1.3),
                    Inches(2), Inches(0.4),
                    m.get("date", ""), size=12, bold=True,
                    color=col, align="center")
        add_textbox(slide, Inches(cx - 1.3), Inches(line_y - 0.9),
                    Inches(2.6), Inches(0.5),
                    m.get("label", ""), size=10,
                    color=p["text"], align="center")
        add_textbox(slide, Inches(cx - 1), Inches(line_y + 0.35),
                    Inches(2), Inches(0.3),
                    (m.get("status", "") or "").upper(), size=8, italic=True,
                    color=p["muted"], align="center")


@recipe("chart_bar")
def chart_bar(slide, params, theme):
    """Native PowerPoint bar/column chart.
    Params: {title, categories[], series: [{name, values[]}], orientation}
    """
    p = _palette(theme)
    title = params.get("title", "")
    categories = params.get("categories", [])
    series = params.get("series", [])

    if title:
        add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.5),
                    title, size=18, bold=True, color=p["primary"])

    if not categories or not series:
        add_textbox(slide, Inches(0.6), Inches(3.5), Inches(12.1), Inches(0.5),
                    "(chart_bar: no data provided)", size=14, color="#999999",
                    align="center")
        return

    data = CategoryChartData()
    data.categories = categories
    for s in series:
        data.add_series(s.get("name", ""), s.get("values", []))

    ctype = (XL_CHART_TYPE.BAR_CLUSTERED if params.get("orientation") == "horizontal"
             else XL_CHART_TYPE.COLUMN_CLUSTERED)
    frame = slide.shapes.add_chart(ctype, Inches(1), Inches(1.8),
                                    Inches(11.3), Inches(5.2), data)
    chart = frame.chart
    chart.has_title = False
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False


@recipe("asks")
def asks(slide, params, theme):
    """Numbered ask cards for executive decks.
    Params: {title, asks: [{number?, title, body, why?, color?}]}
    """
    p = _palette(theme)
    title = params.get("title", "")
    items = params.get("asks", [])

    if title:
        add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.5),
                    title, size=18, bold=True, color=p["primary"])

    n = len(items)
    if n == 0:
        return

    y0 = 1.7
    total_h = 5.4
    gap = 0.2
    card_h = min((total_h - gap * (n - 1)) / n, 2.0)
    palette_cycle = [p["primary"], p["accent"], p["accent2"], p["accent3"]]

    for i, a in enumerate(items):
        y = y0 + i * (card_h + gap)
        color = a.get("color") or palette_cycle[i % len(palette_cycle)]
        add_rect(slide, Inches(0.6), Inches(y), Inches(12.1), Inches(card_h),
                 "#FFFFFF", line_hex="#BFBFBF")
        add_rect(slide, Inches(0.6), Inches(y), Inches(0.9), Inches(card_h), color)
        add_textbox(slide, Inches(0.6), Inches(y), Inches(0.9), Inches(card_h),
                    f"#{a.get('number', i + 1)}", size=32, bold=True,
                    color="#FFFFFF", align="center", anchor="middle")
        add_textbox(slide, Inches(1.7), Inches(y + 0.15),
                    Inches(10.9), Inches(0.45),
                    a.get("title", ""), size=14, bold=True, color=color)
        add_textbox(slide, Inches(1.7), Inches(y + 0.6),
                    Inches(10.9), Inches(card_h - 1.05),
                    a.get("body", ""), size=11, color=p["text"])
        if a.get("why"):
            add_textbox(slide, Inches(1.7), Inches(y + card_h - 0.45),
                        Inches(10.9), Inches(0.4),
                        f"Why: {a['why']}", size=10, italic=True, color=p["muted"])


# ----------------------------- quadrant -----------------------------

@recipe("quadrant")
def quadrant(slide, params, theme):
    """2x2 matrix with labeled axes and 4 cells.

    Params:
      title: str
      x_axis: [left_label, right_label]
      y_axis: [top_label, bottom_label]
      cells: [top_left, top_right, bottom_left, bottom_right]
             each cell is {title, items?} OR a plain string
      highlight: "TL"|"TR"|"BL"|"BR" (optional) — one cell gets the accent fill
    """
    p = _palette(theme)
    title = params.get("title", "")
    x_axis = params.get("x_axis", ["Low", "High"])
    y_axis = params.get("y_axis", ["High", "Low"])
    cells = params.get("cells", [{}, {}, {}, {}])
    highlight = (params.get("highlight") or "").upper()

    if title:
        add_textbox(slide, Inches(0.6), Inches(0.9), Inches(12.1), Inches(0.5),
                    title, size=18, bold=True, color=p["primary"])

    # Grid bounds
    gx, gy = 1.7, 1.6
    gw, gh = 10.8, 5.0
    cw, ch = gw / 2, gh / 2

    fills = {"TL": "#FFFFFF", "TR": "#FFFFFF", "BL": "#FFFFFF", "BR": "#FFFFFF"}
    if highlight in fills:
        fills[highlight] = p["accent3"]

    positions = {
        "TL": (gx, gy),
        "TR": (gx + cw, gy),
        "BL": (gx, gy + ch),
        "BR": (gx + cw, gy + ch),
    }

    for key in ("TL", "TR", "BL", "BR"):
        x, y = positions[key]
        add_rect(slide, Inches(x), Inches(y), Inches(cw), Inches(ch),
                 fills[key], line_hex="#BFBFBF")

    # Cell content
    cell_keys = ["TL", "TR", "BL", "BR"]
    for i, cell in enumerate(cells[:4]):
        key = cell_keys[i]
        x, y = positions[key]
        if isinstance(cell, str):
            text = cell
            items = []
        else:
            text = cell.get("title", "")
            items = cell.get("items", [])
        if text:
            add_textbox(slide, Inches(x + 0.2), Inches(y + 0.2),
                        Inches(cw - 0.4), Inches(0.4),
                        text, size=13, bold=True, color=p["primary"])
        if items:
            tb = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y + 0.7),
                Inches(cw - 0.4), Inches(ch - 0.9))
            tf = tb.text_frame
            tf.word_wrap = True
            for j, item in enumerate(items):
                par = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                par.line_spacing = 1.2
                run = par.add_run()
                run.text = f"*  {item}"
                run.font.size = Pt(11)
                run.font.color.rgb = hex_to_rgb(p["text"])

    # Axis labels
    add_textbox(slide, Inches(gx), Inches(gy + gh + 0.05),
                Inches(cw), Inches(0.3),
                str(x_axis[0]), size=10, italic=True,
                color=p["muted"], align="center")
    add_textbox(slide, Inches(gx + cw), Inches(gy + gh + 0.05),
                Inches(cw), Inches(0.3),
                str(x_axis[1]), size=10, italic=True,
                color=p["muted"], align="center")
    add_textbox(slide, Inches(gx - 1.05), Inches(gy),
                Inches(1.0), Inches(0.3),
                str(y_axis[0]), size=10, italic=True,
                color=p["muted"], align="right")
    add_textbox(slide, Inches(gx - 1.05), Inches(gy + ch),
                Inches(1.0), Inches(0.3),
                str(y_axis[1]), size=10, italic=True,
                color=p["muted"], align="right")


# ----------------------------- process_flow -----------------------------

@recipe("process_flow")
def process_flow(slide, params, theme):
    """Horizontal arrow chain of phases.

    Params:
      title: str
      phases: list of {name, description?, color?}
    """
    p = _palette(theme)
    title = params.get("title", "")
    phases = params.get("phases", [])

    if title:
        add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.5),
                    title, size=18, bold=True, color=p["primary"])

    if not phases:
        return

    from pptx.enum.shapes import MSO_SHAPE
    n = len(phases)
    total_w = 12.1
    gap = 0.15
    phase_w = (total_w - gap * (n - 1)) / n
    y = 2.5
    h = 2.5
    palette_cycle = [p["primary"], p["accent"], p["accent2"], p["accent3"]]

    for i, phase in enumerate(phases):
        x = 0.6 + i * (phase_w + gap)
        color = phase.get("color") or palette_cycle[i % len(palette_cycle)]

        # Use pentagon (right arrow shape) except last — use rounded rect for end
        shape_type = MSO_SHAPE.PENTAGON if i < n - 1 else MSO_SHAPE.ROUNDED_RECTANGLE
        arrow = slide.shapes.add_shape(
            shape_type,
            Inches(x), Inches(y), Inches(phase_w), Inches(h))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = hex_to_rgb(color)
        arrow.line.fill.background()

        # Phase number
        add_textbox(slide, Inches(x + 0.2), Inches(y + 0.25),
                    Inches(phase_w - 0.4), Inches(0.4),
                    f"PHASE {i + 1}", size=9, bold=True,
                    color="#FFFFFF", align="center")
        # Name
        add_textbox(slide, Inches(x + 0.2), Inches(y + 0.7),
                    Inches(phase_w - 0.4), Inches(0.5),
                    phase.get("name", ""), size=16, bold=True,
                    color="#FFFFFF", align="center", anchor="middle")
        # Description
        if phase.get("description"):
            add_textbox(slide, Inches(x + 0.3), Inches(y + 1.35),
                        Inches(phase_w - 0.6), Inches(h - 1.45),
                        phase["description"], size=10,
                        color="#FFFFFF", align="center")


# ----------------------------- section_divider -----------------------------

@recipe("section_divider")
def section_divider(slide, params, theme):
    """Full-bleed section header.

    Params:
      number: str or int (e.g., "01", "II", "3")
      title: str
      subtitle: str (optional)
      color: hex (optional — overrides accent)
    """
    p = _palette(theme)
    color = params.get("color") or p["primary"]
    number = str(params.get("number", ""))
    title = params.get("title", "")
    subtitle = params.get("subtitle", "")

    # Full-bleed color panel
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), color)

    # Left accent stripe
    add_rect(slide, Inches(0), Inches(0), Inches(0.5), Inches(7.5), p["accent3"])

    # Number (large, faded)
    if number:
        add_textbox(slide, Inches(1.0), Inches(0.8), Inches(4), Inches(2),
                    number, size=72, bold=True, color="#FFFFFF",
                    anchor="top")

    # Title
    add_textbox(slide, Inches(1.0), Inches(3.0), Inches(11), Inches(2),
                title, size=48, bold=True, color="#FFFFFF")

    # Subtitle
    if subtitle:
        add_textbox(slide, Inches(1.0), Inches(5.2), Inches(11), Inches(1),
                    subtitle, size=20, italic=True, color="#FFFFFF")
