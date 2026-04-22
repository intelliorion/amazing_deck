"""Deck generation - content JSON + (optional) template -> .pptx."""
from pathlib import Path
from pptx import Presentation

from .render import fill_placeholder, apply_extra
from .recipes import apply_recipe
from .extract import extract_theme


def generate_deck(template_path, content, output_path, verbose=True,
                  keep_template_slides=False):
    if template_path:
        prs = Presentation(str(template_path))
        theme = extract_theme(prs)
        if not keep_template_slides:
            removed = _clear_existing_slides(prs)
            if verbose and removed:
                print(f"Cleared {removed} pre-existing slide(s) from template "
                      f"(use --keep-template-slides to preserve)")
    else:
        prs = Presentation()
        prs.slide_width = 12192000
        prs.slide_height = 6858000
        theme = {"colors": {}, "fonts": {}}

    blank_idx = _find_blank(prs)
    slides_spec = content.get("slides", [])

    for i, spec in enumerate(slides_spec):
        if "recipe" in spec:
            layout = prs.slide_layouts[blank_idx]
            slide = prs.slides.add_slide(layout)
            apply_recipe(slide, spec["recipe"],
                         spec.get("recipe_params", {}), theme)
            label = f"recipe:{spec['recipe']}"
        else:
            layout_ref = spec.get("layout", 0)
            layout = _resolve_layout(prs, layout_ref)
            slide = prs.slides.add_slide(layout)
            for k, v in spec.get("placeholders", {}).items():
                ok = fill_placeholder(slide, k, v)
                if not ok and verbose:
                    print(f"  [slide {i+1}] warn: placeholder '{k}' not matched")
            label = f"layout:{layout.name}"

        for extra in spec.get("extras", []):
            apply_extra(slide, extra)

        if verbose:
            print(f"  Slide {i+1}: {label}")

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    if verbose:
        print(f"\nSaved: {output_path}   ({len(prs.slides)} slides)")
    return output_path


def _resolve_layout(prs, ref):
    """Find a layout by index, exact name, or fuzzy match.

    Matching order:
      1. Numeric index
      2. Exact name (case-insensitive)
      3. Substring match either direction ("Title Slide" finds "Title")
      4. Word-overlap best match ("Title and Content" finds "Title and Content 1")
    """
    if isinstance(ref, int):
        return prs.slide_layouts[ref]

    ref_lower = str(ref).lower().strip()
    layouts = list(prs.slide_layouts)

    # 1. Exact
    for layout in layouts:
        if layout.name.lower() == ref_lower:
            return layout

    # 2. Substring match (prefer shorter = more specific)
    subs = []
    for layout in layouts:
        ln = layout.name.lower()
        if ref_lower in ln or ln in ref_lower:
            subs.append(layout)
    if subs:
        subs.sort(key=lambda x: abs(len(x.name) - len(ref_lower)))
        return subs[0]

    # 3. Word-overlap
    def words(s):
        return set(s.lower().replace("+", " ").replace("-", " ").replace("_", " ").split())

    ref_words = words(ref_lower)
    best = None
    best_score = 0
    for layout in layouts:
        overlap = len(ref_words & words(layout.name))
        if overlap > best_score:
            best_score = overlap
            best = layout

    if best and best_score > 0:
        print(f"  [info] layout '{ref}' -> '{best.name}' (fuzzy match, {best_score} word(s))")
        return best

    print(f"  [warn] layout '{ref}' not found, using layout[0] '{layouts[0].name}'")
    return layouts[0]


def _clear_existing_slides(prs):
    """Remove all pre-existing slides from a template. Returns count removed.

    Properly drops both the sldId element AND the relationship to the slide
    part, so the underlying XML parts are garbage-collected on save (avoids
    duplicate-name warnings from the zip writer).
    """
    from pptx.oxml.ns import qn
    xml_slides = prs.slides._sldIdLst
    existing = list(xml_slides)
    for sld_id in existing:
        r_id = sld_id.get(qn("r:id"))
        if r_id:
            try:
                prs.part.drop_rel(r_id)
            except Exception:
                pass
        xml_slides.remove(sld_id)
    return len(existing)


def _find_blank(prs):
    for i, layout in enumerate(prs.slide_layouts):
        if "blank" in layout.name.lower():
            return i
    return len(prs.slide_layouts) - 1
