"""Placeholder filling + basic shape rendering."""
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from .utils import hex_to_rgb


ROLE_TO_TYPES = {
    "title": ("TITLE", "CENTER_TITLE", "CTR_TITLE"),
    "subtitle": ("SUBTITLE",),
    "content": ("BODY", "OBJECT"),
    "body": ("BODY", "OBJECT"),
    "date": ("DATE",),
    "footer": ("FOOTER",),
    "slide_number": ("SLIDE_NUMBER",),
}


def _norm_type(t):
    """Normalize PP_PLACEHOLDER enum string across python-pptx versions."""
    if t is None:
        return ""
    try:
        name = getattr(t, "name", None)
        if name:
            return name
    except Exception:
        pass
    s = str(t)
    s = s.split(" ", 1)[0]
    s = s.rsplit(".", 1)[-1]
    return s


def fill_placeholder(slide, key, value):
    target = find_placeholder(slide, key)
    if target is None:
        return False

    if isinstance(value, list):
        _fill_bullets(target, value)
    elif isinstance(value, dict):
        _fill_rich(slide, target, value)
    else:
        target.text = str(value)
    return True


def find_placeholder(slide, key):
    key_str = str(key).lower()
    for ph in slide.placeholders:
        ph_type = _norm_type(ph.placeholder_format.type)
        for role, types in ROLE_TO_TYPES.items():
            if key_str == role and ph_type in types:
                return ph
    for ph in slide.placeholders:
        if ph.name.lower() == key_str:
            return ph
    try:
        idx = int(key)
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                return ph
    except (ValueError, TypeError):
        pass
    return None


def _fill_bullets(placeholder, items):
    tf = placeholder.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(item)


def _fill_rich(slide, placeholder, value):
    t = value.get("type")
    if t == "table":
        add_table(slide, placeholder.left, placeholder.top,
                  placeholder.width, placeholder.height, value["rows"])
        if placeholder.has_text_frame:
            placeholder.text_frame.clear()
    elif t == "image":
        slide.shapes.add_picture(value["path"], placeholder.left, placeholder.top,
                                  placeholder.width, placeholder.height)
        if placeholder.has_text_frame:
            placeholder.text_frame.clear()


def add_table(slide, x, y, w, h, rows, header_color="#1B3E6F",
              header_text="#FFFFFF"):
    if not rows:
        return None
    n_rows, n_cols = len(rows), len(rows[0])
    shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = shape.table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb(header_color)
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.bold = True
                        run.font.color.rgb = hex_to_rgb(header_text)
    return tbl


def add_textbox(slide, x, y, w, h, text, *, size=14, bold=False,
                color="#000000", align="left", anchor="top", italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    anchors = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
               "bottom": MSO_ANCHOR.BOTTOM}
    aligns = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
              "right": PP_ALIGN.RIGHT}
    tf.vertical_anchor = anchors.get(anchor, MSO_ANCHOR.TOP)
    p = tf.paragraphs[0]
    p.alignment = aligns.get(align, PP_ALIGN.LEFT)
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = hex_to_rgb(color)
    return tb


def add_rect(slide, x, y, w, h, fill_hex, line_hex=None):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    if line_hex is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = hex_to_rgb(line_hex)
    s.shadow.inherit = False
    return s


def apply_extra(slide, extra):
    t = extra.get("type")
    x = Inches(extra["x_in"])
    y = Inches(extra["y_in"])
    w = Inches(extra["w_in"])
    h = Inches(extra["h_in"])
    if t == "textbox":
        return add_textbox(slide, x, y, w, h, extra.get("text", ""),
                           size=extra.get("font_size", 14),
                           bold=extra.get("bold", False),
                           italic=extra.get("italic", False),
                           color=extra.get("color", "#000000"),
                           align=extra.get("align", "left"),
                           anchor=extra.get("anchor", "top"))
    if t == "rectangle":
        return add_rect(slide, x, y, w, h, extra.get("fill", "#FFFFFF"),
                        line_hex=extra.get("line"))
    if t == "table":
        return add_table(slide, x, y, w, h, extra["rows"],
                         header_color=extra.get("header_color", "#1B3E6F"))
    return None
